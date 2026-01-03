import numpy as np
import pandas as pd
from bigquery import get_bigquery_client
from config import REVERSE_METRICS, config_dataset, config_info_table, table_suffix, log_table
from pptx.util import Inches, Cm, Pt
from pptx.dml.color import RGBColor
from anytree import AbstractStyle
from pptx.enum.text import MSO_ANCHOR
from datetime import timedelta

from google.cloud import bigquery

class DoubleStyle(AbstractStyle):

    def __init__(self):
        super(DoubleStyle, self).__init__(u'       \u2551   ',
                                          u'       \u2560\u2550\u2550 ',
                                          u'       \u255a\u2550\u2550 ')


def delta_pct(now, prev):
    if np.isnan(prev):
        prev = 0

    if np.isnan(now):
        now = 0

    if prev == 0:
        if now > 0:
            return float('inf')
        elif now < 0:
            return -float('inf')
        else:
            return 0
    else:
        return (now / prev - 1) * 100


def delta_print_pct(prev, now):
    prev = float(prev)
    now = float(now)
    if prev == 0:
        if now == 0:
            return 0
        else:
            return "∞"
    else:
        result = now / prev - 1
        if abs(result) < 10:
            return abs(round((now / prev - 1) * 100, 1))
        else:
            return ">1000"


def print_delta(prev, now):
    greater = now > prev
    return f"{'▲' if greater else '▼'}{delta_print_pct(prev=prev, now=now)}%"
    # return f"{'+' if greater else '-'}{delta_print_pct(prev=prev, now=now)}%"


def human_format(num):
    if num == np.inf:
        return '∞'

    if num == -np.inf:
        return '-∞'

    sign = '-' if num < 0 else ''

    num = abs(num)
    magnitude = 0
    original_num = num
    while abs(num) >= 1000:
        magnitude += 1
        num /= 1000.0
    # add more suffixes if you need them
    if (magnitude == 1):
        return f"{sign}{(round(num, 2))}K"
    elif (magnitude == 2):
        return f"{sign}{round(num, 2)}M"
    else:
        return f"{sign}{int(original_num)}"


def fix_name(metric):
    if metric.lower() == 'ctr':
        return 'Click Through Rate'

    if metric.lower() == 'cpc':
        return 'Cost Per Click'

    if metric.lower() == 'clicks':
        return 'Traffic'

    if metric.lower() == 'orders':
        return 'Total-Orders'

    if metric.lower() == 'aov':
        return 'Total-AOV'

    if '__to__' in metric:
        metric = ' '.join(metric.split('_'))
        result = []
        for word in metric.split():
            result.append(word.title())
        return ' '.join(result) + ' CVR'
    else:
        return '-'.join(metric.split('_'))


def get_metric_type(metric):
    metric = metric.lower()

    if metric == 'revenue_share':
        return 'rate'

    if metric.startswith('revenue') or metric.endswith('revenue') or metric.endswith('spend') or metric.endswith(
            'cost') or metric.endswith('returns') or metric.startswith('Revenue'):
        return 'revenue'
    elif metric.startswith('aov') or metric.startswith('cpc') or metric.startswith('cac'):
        return 'aov'
    elif metric.endswith('rate') or ('__to__' in metric) or metric == 'ctr' or metric.endswith(
            'share') or metric == 'acos':
        return 'rate'
    elif metric.endswith('duration'):
        return 'time'
    else:
        return 'traffic'


def print_formatted(value, metric, currency=None):
    metric = metric.lower()

    if np.isnan(value):
        return ""

    if currency is None:
        currency = "$"

    metric_type = get_metric_type(metric)

    if metric in ['upt', 'aur'] or 'aov' in metric:
        if 'aov' in metric:
            return f"{currency}{human_format(value)}"
        return f"{value:.2f}"
    if metric_type == 'revenue':
        return f"{currency}{human_format(value)}"
    elif metric_type == 'aov':
        return f"{currency}{value:.2f}"
    elif metric_type == 'rate':
        # if value <= 1:
        #     return f"{100*value:.2f}%"
        # else:
        #     return f"{value:.2f}%"
        return f"{100*value:.2f}%"
    elif metric_type == 'time':
        time_str = str(timedelta(seconds=value))
        # 0 hours
        if time_str[0] == '0':
            return f"{time_str[2:7]}"
        else:
            return f"{time_str[:7]}"
    else:
        return f"{human_format(value)}"


def divide(a, b):
    if b != 0:
        return a/b
    else:
        return 0


def filter_dims(kpi_df, dim_values_to_include):
    for dim in dim_values_to_include:
        mask = kpi_df[dim].isin(dim_values_to_include[dim])
        kpi_df = kpi_df[mask]

    return kpi_df


def not_none(number):
    return (number == number) and number is not None


def is_none(number):
    return not not_none(number)


def get_tables(project_id, dataset_id, period):
    result = []
    client = get_bigquery_client(project_id)
    table_ids = [table.table_id for table in client.list_tables(dataset_id)]
    for table_id in table_ids:
        if table_id.endswith('_view') and (period in table_id) and ('raw_funnel' not in table_id):
            result.append(table_id[:-5] + table_suffix)

    return result


def get_dim_metrics(project_id, dataset_id, table_id):
    client = get_bigquery_client(project_id)
    dataset_ref = client.dataset(dataset_id, project=project_id)
    table_ref = dataset_ref.table(table_id)
    table = client.get_table(table_ref)

    dims = []
    metrics = []
    for schema in table.schema:
        if schema.field_type in ['NUMERIC', 'FLOAT', 'INTEGER']:
            if schema.name.endswith('_yhat'):
                metrics.append(schema.name[:-5])
        elif schema.name.lower() in ['date', 'datehour', 'week']:
            continue
        else:
            dims.append(schema.name)

    if len(dims) == 0:
        return None, metrics

    elif len(dims) == 1:
        return dims[0], metrics

    else:
        raise ValueError(f'Error while getting dim and metrics for {project_id}.{dataset_id}.{table_id} : length of dims is {len(dims)}')


def get_anomaly_df(project_id, dataset_id, table_id, period, date_filter=None):
    if period == 'hourly':
        date_col = 'DateHour'
    elif period == 'daily':
        date_col = 'Date'
        if 'return' in table_id:
            data_col = 'Return_Date'
    elif period == 'weekly':
        date_col = 'Week'
    else:
        raise Exception(f"Invalid period - {period}")

    if date_filter:
        query = f"""
                SELECT * FROM
                `{project_id}.{dataset_id}.{table_id}`
                WHERE {date_col} {date_filter}
                ORDER BY {date_col}
                """
    else:
        query = f"""
                SELECT * FROM
                `{project_id}.{dataset_id}.{table_id}`
                ORDER BY {date_col}
                """

    client = get_bigquery_client(project_id)
    anomaly_df = (
        client.query(query)
            .result()
            .to_dataframe()
    )

    _, metrics = get_dim_metrics(project_id, dataset_id, table_id)

    for metric in metrics:
        anomaly_df[metric] = pd.to_numeric(anomaly_df[metric])

    return anomaly_df


def get_table_details(project_id, asset, table_id):
    data_source = get_data_source(table_id)
    dim, metrics = get_dim_metrics(project_id, asset, table_id)

    return data_source, dim, metrics


def get_anomaly_type(y, upper, lower):
    if y is None:
        return 0

    if y > upper:
        return 1
    elif y < lower:
        return -1
    else:
        return 0


def check_warning(y, upper, lower, lower_bound=10, upper_bound=30):
    return (lower_bound < -delta_pct(y, lower) <= upper_bound) or (lower_bound < delta_pct(y, upper) <= upper_bound)


def check_critical(y, upper, lower, threshold=30):
    return (-delta_pct(y, lower) > threshold) or (delta_pct(y, upper) > threshold)


def get_color(anomaly_type, metric):
    if metric.lower() in REVERSE_METRICS:
        anomaly_type *= -1

    if anomaly_type == 1:
        return 'green'
    elif anomaly_type == -1:
        return 'red'
    else:
        return None


def get_data_source(table_id):
    if table_id.startswith('ga_'):
        data_source = 'Google Analytics'
    elif table_id.startswith('fb_'):
        data_source = 'Facebook'
    elif table_id.startswith('googleAds_'):
        data_source = 'Google Ads'
    elif table_id.startswith('custom_') or table_id.startswith('shopify_') or table_id.startswith('ecommerce_'):
            data_source = 'Ecommerce'
    elif table_id.startswith('magento_'):
            data_source = 'Magento'
    elif table_id.startswith('upscribe_'):
        data_source = 'Upscribe'
    elif table_id.startswith('affiliate_'):
        data_source = 'Affiliate'
    else:
        data_source = table_id.split('_')[0]

    return data_source


def add_color(txt, color):
    if color:
        # return f"""<span style="color:{color}">{txt}</span>"""
        return f"""<font color="{color}">{txt}</font>"""
    else:
        return txt


class Element:
    def __init__(self, text, data_type='regular_text'):
        self.text = text
        self.type = data_type


def bold(txt):
    return f'<b>{txt}</b>'


def main_heading(text):
    text = bold(text)
    # return f'<h1 style="font-size:25px">{text}</h1><br>'
    # return f"<h1>{add_color(text, '#2980b9')}</h1><br>"
    return f"<h1>{add_color(text, '#2980b9')}</h1>"


def sub_heading(text):
    text = bold(text)
    # return f'<br><h1 style="font-size:23px">{text}</h1><br>'
    # return f"<h3>{add_color(text, '#2980b9')}</h3><br>"
    return f"<h3>{add_color(text, '#2980b9')}</h3>"


def regular_text(text):
    return f'<p>{text}</p>'


always_include_data_source = ['Google Ads', 'Facebook']


def print_anomaly(p, row, pre):

    data_source = row.data_source
    dimension = row.dimension
    dim_label = row.dim_label
    metric = row.metric
    period = row.period
    is_warning = row.is_warning
    is_critical = row.is_critical
    y = row.y
    y_prev = row.y_prev
    yhat = row.yhat
    yhat_upper = row.yhat_upper
    yhat_lower = row.yhat_lower
    revenue_impact = row.revenue_impact
    reverse_effect_on_parent = row.reverse_effect_on_parent
    fact_vs_forecast = ''

    prev_anomaly_type = row.anomaly_type
    prev_color = get_color(prev_anomaly_type, metric)

    yhat_anomaly_type = get_anomaly_type(y=y, upper=yhat_upper, lower=yhat_lower)
    yhat_color = get_color(yhat_anomaly_type, metric)

    if data_source in always_include_data_source:
        fact_vs_forecast = fact_vs_forecast + f'{data_source} '

    if not_none(dim_label):
        # fact_vs_forecast = fact_vs_forecast + bold(f'"{fix_name(dimension)}: {fix_name(dim_label)}"')
        fact_vs_forecast = fact_vs_forecast + f'{fix_name(dimension)} - {fix_name(dim_label)}'
    else:
        fact_vs_forecast = fact_vs_forecast + f" {fix_name(metric)} {print_formatted(y, metric)}"

    run = p.add_run()
    run.text = '\n' + f'{pre}{fact_vs_forecast}'

    font = run.font
    font.name = 'Poppins'
    font.size = Pt(10)
    font.bold = True
    font.color.rgb = RGBColor(100, 100, 100)

    if yhat_color == 'red':
        run = p.add_run()
        # run.text = f" ({print_delta(now=y, prev=yhat)}) - Revenue Impact : " \
        #            f"{print_formatted(revenue_impact, 'Total_Sales')}"

        run.text = f" ({print_delta(now=y, prev=yhat)})"

        font = run.font
        font.name = 'Poppins'
        font.size = Pt(10)
        font.bold = True
        font.color.rgb = RGBColor(229, 88, 1) if is_warning else RGBColor(192, 0, 0)
    elif yhat_color == 'green':
        run = p.add_run()
        # run.text = f" ({print_delta(now=y, prev=yhat)}) - Revenue Impact : " \
        #            f"{print_formatted(revenue_impact, 'Total_Sales')}"

        run.text = f" ({print_delta(now=y, prev=yhat)})"

        font = run.font
        font.name = 'Poppins'
        font.size = Pt(10)
        font.bold = True
        if reverse_effect_on_parent:
            font.color.rgb = RGBColor(0, 220, 0)
        else:
            font.color.rgb = RGBColor(0, 200, 0)
    else:
        run = p.add_run()
        # run.text = f" ({print_delta(now=y, prev=yhat)}) - Revenue Impact : " \
        #            f"{print_formatted(revenue_impact, 'Total_Sales')}"

        run.text = f" ({print_delta(now=y, prev=yhat)})"

        font = run.font
        font.name = 'Poppins'
        font.size = Pt(10)
        font.bold = True
        # if yhat_lower > y:
        #     font.color.rgb = RGBColor(0, 150, 0)
        # else:
        #     font.color.rgb = RGBColor(150, 0, 0)
        font.color.rgb = RGBColor(172, 172, 172)


def print_comment(row):
    dimension = row.dimension
    dim_label = row.dim_label
    metric = row.metric
    period = row.period
    is_warning = row.is_warning
    is_critical = row.is_critical
    y = row.y
    y_prev = row.y_prev
    yhat = row.yhat
    yhat_upper = row.yhat_upper
    yhat_lower = row.yhat_lower

    if period == 'daily':
        comparison = 'SDLW'
    elif period == 'weekly':
        comparison = 'WoW'
    else:
        raise Exception(f"Invalid period - {period}")

    comment = ''

    if not_none(dim_label):
        # fact_vs_forecast = fact_vs_forecast + bold(f'"{fix_name(dimension)}: {fix_name(dim_label)}"')
        comment = comment + f'"{fix_name(dim_label)}"'
    comment = comment + f" {fix_name(metric)}"

    if y > y_prev:
        comment = comment + f" increased by "
    else:
        comment = comment + f" decreased by "

    comment = comment + f"{print_formatted(y-y_prev, metric)}"

    comment = comment + f" {comparison}"

    return comment


def print_orders(client_name, orders_df, asset_df, p, orders_type, sales=None):

    style = DoubleStyle()
    client_orders_df = orders_df[orders_df.client_name == client_name]
    total_orders = float(client_orders_df['Total_Orders'].values[0])
    organic_orders = float(client_orders_df['Organic_Orders'].values[0])
    ad_orders = float(client_orders_df['Ad_Orders'].values[0])

    if orders_type == 'Organic':
        run = p.add_run()
        run.text = f"\n{style.vertical}{style.cont} Organic-AOV {print_formatted(sales/organic_orders, 'AOV')}"
        font = run.font
        font.name = 'Poppins'
        font.size = Pt(10)
        font.bold = True
        font.color.rgb = RGBColor(100, 100, 100)
        run = p.add_run()
        run.text = f"\n{style.vertical}{style.end} Organic-Orders {print_formatted(organic_orders, 'Orders')}"
        font = run.font
        font.name = 'Poppins'
        font.size = Pt(10)
        font.bold = True
        font.color.rgb = RGBColor(100, 100, 100)
        run = p.add_run()
        run.text = f" ({print_formatted(organic_orders/total_orders, 'rate')})"
        font = run.font
        font.name = 'Poppins'
        font.size = Pt(10)
        font.bold = True
        font.color.rgb = RGBColor(172, 172, 172)

    else:
        if sales is None:
            ad_sales = asset_df[(asset_df.metric == 'Ad_Sales') & (asset_df.dimension.isnull())].y.values[0]
        else:
            ad_sales = sales
        run = p.add_run()
        run.text = f"\n{style.empty}{style.cont} Ad-AOV {print_formatted(ad_sales / ad_orders, 'AOV')}"
        font = run.font
        font.name = 'Poppins'
        font.size = Pt(10)
        font.bold = True
        font.color.rgb = RGBColor(100, 100, 100)
        run = p.add_run()
        run.text = f"\n{style.empty}{style.end} Ad-Orders {print_formatted(ad_orders, 'Orders')}"
        font = run.font
        font.name = 'Poppins'
        font.size = Pt(10)
        font.bold = True
        font.color.rgb = RGBColor(100, 100, 100)
        run = p.add_run()
        run.text = f" ({print_formatted(ad_orders / total_orders, 'rate')})"
        font = run.font
        font.name = 'Poppins'
        font.size = Pt(10)
        font.bold = True
        font.color.rgb = RGBColor(172, 172, 172)


def print_sales(client_name, asset_df, orders_df, p):
    total_sales = asset_df[(asset_df.metric == 'Total_Sales') & (asset_df.dimension.isnull())].y.values[0]
    ad_sales = asset_df[(asset_df.metric == 'Ad_Sales') & (asset_df.dimension.isnull())].y.values[0]
    organic_sales = total_sales - ad_sales

    style = DoubleStyle()
    run = p.add_run()
    run.text = f"\n{style.cont} Organic-Sales {print_formatted(organic_sales, 'Orders')}"
    font = run.font
    font.name = 'Poppins'
    font.size = Pt(10)
    font.bold = True
    font.color.rgb = RGBColor(100, 100, 100)
    # run = p.add_run()
    # run.text = f" ({print_formatted(organic_sales / total_sales, 'rate')})"
    # font = run.font
    # font.name = 'Poppins'
    # font.size = Pt(12)
    # font.bold = True
    # font.color.rgb = RGBColor(172, 172, 172)
    if client_name in orders_df.client_name.unique():
        print_orders(client_name, orders_df, asset_df, p, orders_type="Organic", sales=organic_sales)
    # run = p.add_run()
    # run.text = f" ({print_formatted(organic_orders / total_orders, 'rate')})"
    # font = run.font
    # font.name = 'Poppins'
    # font.size = Pt(12)
    # font.bold = True
    # font.color.rgb = RGBColor(172, 172, 172)


def filter_data_by_kpi(asset_df, kpi):
    data_source_mask = asset_df['data_source'] == kpi.data_source

    if kpi.dimension:
        dimension_mask = asset_df['dimension'] == kpi.dimension
    else:
        dimension_mask = asset_df['dimension'].isnull()

    if kpi.dim_label:
        dim_label_mask = asset_df['dim_label'] == kpi.dim_label
    else:
        dim_label_mask = asset_df['dim_label'].isnull()

    metric_mask = asset_df['metric'] == kpi.metric

    filtered_df = asset_df[data_source_mask & dimension_mask & dim_label_mask & metric_mask]
    return filtered_df


def add_anomaly_legend(anomaly_slide):

    legend_shape_left = Inches(9.3)
    legend_shape_top = Inches(6.4)
    legend_shape_width = Cm(4.81)
    legend_shape_height = Cm(6.98)

    legend_textbox = anomaly_slide.shapes.add_textbox(
        legend_shape_left, legend_shape_top, legend_shape_width, legend_shape_height
    )

    legend_text_frame = legend_textbox.text_frame
    legend_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = legend_text_frame.paragraphs[0]

    list_shape = ['▼', '▼', '▲', '▲', '▲', '▼']
    list_comments = ['-Critical anomaly alert that requires immediate attention.',
                     '-Warning anomaly alert with lower negative revenue impact.',
                     '-Anomaly with large positive revenue impact',
                     '-Anomaly with moderate positive revenue impact',
                     '-Increase within expected range',
                     '-Decrease within expected range']
    list_colors = [RGBColor(192, 0, 0), RGBColor(229, 88, 1), RGBColor(0, 200, 0), RGBColor(0, 220, 0),
                   RGBColor(172, 172, 172), RGBColor(172, 172, 172)]

    for shape, comment, color in zip(list_shape, list_comments, list_colors):

        legend_text_frame.add_paragraph()
        para = legend_text_frame.paragraphs[list_comments.index(comment)]
        run = para.add_run()

        run.text = shape
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(9)
        font.bold = True
        font.color.rgb = color

        run = para.add_run()
        run.text = comment
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(9)
        font.bold = True

def update_story_sent_info(project_id, dataset_id, timezone, status):
    
    try:
        client = get_bigquery_client(project_id)
        query_string = f"""
            UPDATE `{project_id}.{config_dataset}.{config_info_table}`
            SET story_sent = {status},
            story_sent_time = TIMESTAMP(DATETIME(CURRENT_TIMESTAMP(),'{timezone}'))
            WHERE view_dataset_id = '{dataset_id}'
            """
        client.query(query_string).result()
    except:
        print("update_story_sent_info Fn. failed to execute")

def execution_logger(client, project_id,df_data):
    try:
        job_config = bigquery.LoadJobConfig(write_disposition="WRITE_APPEND")
        execution_table_id = f'{project_id}.{config_dataset}.{log_table}'
        job = client.load_table_from_dataframe(df_data, execution_table_id, job_config=job_config)
        print("Execution Log updated", job)
    except Exception as err:
        print("Logger - Writing to BQ Failed",err)
    return
