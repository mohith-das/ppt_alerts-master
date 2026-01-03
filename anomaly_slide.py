import numpy as np
from config import in_production, currency_dict, color_dict, REVERSE_METRICS
from helper import print_formatted, print_delta, fix_name, always_include_data_source, not_none, filter_data_by_kpi, \
add_anomaly_legend, is_none
from dates import yesterday, current_date, get_previous_week_start_date_end_date
from graph import get_graph
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from rca import print_revenue_rca
from collections import OrderedDict


def is_kpi(row, asset_df, kpi_list):
    for kpi in kpi_list:
        filtered_row = filter_data_by_kpi(asset_df, kpi)
        if (not filtered_row.empty) and row.equals(filtered_row.iloc[0, :]):
            return True

    return False


def add_anomaly_heading(slide, asset):
    main_heading_text_left = Inches(0.15)
    main_heading_text_top = Inches(0.05)
    main_heading_text_width = Inches(10)
    main_heading_text_height = Inches(1)
    main_heading_textbox = slide.shapes.add_textbox(
        main_heading_text_left, main_heading_text_top, main_heading_text_width, main_heading_text_height
    )

    main_heading_text_frame = main_heading_textbox.text_frame
    main_heading_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = main_heading_text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f'{fix_name(asset)}-Anomaly Alerts'
    font = run.font
    font.name = 'Poppins'
    font.size = Pt(28)
    font.bold = True
    font.color.rgb = RGBColor(70, 177, 255)


def add_date(slide, period, weekday=None):
    date_text_left = Inches(10.5)
    date_text_top = Inches(0.25)
    date_text_width = Cm(5.89)
    date_text_height = Cm(1.66)
    date_textbox = slide.shapes.add_textbox(
        date_text_left, date_text_top, date_text_width, date_text_height
    )

    date_text_frame = date_textbox.text_frame
    date_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = date_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()

    if period == 'daily':
        date_text = f'Date - {yesterday.strftime("%d/%b/%y")}'
    elif period == 'weekly':
        week_start, week_end = get_previous_week_start_date_end_date(current_date, weekday=weekday)
        date_text = f'Week {week_start.strftime("%U")} - {week_start.strftime("%d/%b/%y")} – {week_end.strftime("%d/%b/%y")}'
    else:
        date_text = f'Unknown period - {period} '

    run.text = date_text
    font = run.font
    font.name = 'Poppins'
    font.size = Pt(14)
    font.bold = True


def add_anomaly_subheadings(anomaly_slide, anomaly_count):
    negative_text_left = Inches(0.15)
    negative_text_top = Inches(0.5)
    negative_text_width = Cm(10.89)
    negative_text_height = Cm(1.66)
    negative_textbox = anomaly_slide.shapes.add_textbox(
        negative_text_left, negative_text_top, negative_text_width, negative_text_height
    )

    text_frame = negative_textbox.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = text_frame.paragraphs[0]
    run = p.add_run()

    run.text = f'Top {anomaly_count} Metrics with most impact on business'
    font = run.font
    font.name = 'Poppins'
    font.size = Pt(14)
    font.bold = True
    font.color.rgb = RGBColor(70, 177, 220)


def add_card_shape(anomaly_slide, left, top, width, height):
    card_shape = anomaly_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    card_shape.adjustments[0] = 0.05

    fill = card_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    card_shape.shadow.inherit = True

    card_shape.line.fill.background()


def add_rca_legend_shape(anomaly_slide, left, top, width, height):
    card_shape = anomaly_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    card_shape.adjustments[0] = 0.05

    fill = card_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    card_shape.shadow.inherit = True

    card_shape.line.fill.background()


def add_rca_legend(anomaly_slide):

    rca_shape_left = Inches(9.25)
    rca_shape_top = Inches(1)
    rca_shape_width = Cm(6.81)
    rca_shape_height = Cm(7.98)


def add_warning_critical(anomaly_slide, left, top, width, height, is_warning, is_critical):
    if not any([is_warning, is_critical]):
        return

    warning_critical_textbox = anomaly_slide.shapes.add_textbox(
        left, top, width, height
    )

    warning_critical_text_frame = warning_critical_textbox.text_frame
    warning_critical_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = warning_critical_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()

    run.text = 'Warning' if is_warning else 'Critical'
    font = run.font
    font.name = 'Poppins'
    font.size = Pt(12)
    font.bold = True
    font.color.rgb = RGBColor(229, 88, 1) if is_warning else RGBColor(192, 0, 0)


def add_metric_name(anomaly_slide, left, top, width, height, metric):
    metric_name_textbox = anomaly_slide.shapes.add_textbox(
        left, top, width, height
    )

    metric_name_text_frame = metric_name_textbox.text_frame
    metric_name_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = metric_name_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()

    metric_name_text = fix_name(metric)

    run.text = metric_name_text
    # if len(run.text) > 22:
    #     run.text = run.text[:19] + '...'
    font = run.font
    font.name = 'Poppins'

    font_size_config = OrderedDict({16: Pt(12), 18: Pt(11), 20: Pt(10)})
    for k, v in font_size_config.items():
        if len(run.text) < k:
            font.size = v
            break
    font.size = Pt(12) if len(run.text) < 20 else Pt(10)
    font.bold = True
    font.color.rgb = RGBColor(166, 166, 166)


def add_metric_value(anomaly_slide, left, top, width, height, y, metric, currency=None):
    metric_value_textbox = anomaly_slide.shapes.add_textbox(
        left, top, width, height
    )

    metric_value_text_frame = metric_value_textbox.text_frame
    metric_value_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = metric_value_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()

    metric_value_text = print_formatted(y, metric, currency)

    run.text = metric_value_text
    font = run.font
    font.name = 'Poppins'
    font_size_config = OrderedDict({8: Pt(16), 10: Pt(14), 12: Pt(12)})
    for k, v in font_size_config.items():
        if len(run.text) < k:
            font.size = v
            break
    font.bold = True


def add_delta(anomaly_slide, left, top, width, height, period, y, yhat, color):
    delta_text = f"{print_delta(now=y, prev=yhat)}"
    if len(delta_text) >= 7:
        left = left - Cm(0.2)
    delta_textbox = anomaly_slide.shapes.add_textbox(
        left, top, width, height
    )

    delta_text_frame = delta_textbox.text_frame
    delta_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = delta_text_frame.paragraphs[0]
    run = p.add_run()

    run.text = delta_text
    font = run.font
    font.name = 'Arial'
    font.size = Pt(11)
    font.bold = True
    font.color.rgb = color_dict[color]


def add_dimension(anomaly_slide, left, top, width, height, dimension, dim_label, data_source):
    dimension_textbox = anomaly_slide.shapes.add_textbox(
        left, top, width, height
    )

    dimension_text_frame = dimension_textbox.text_frame
    dimension_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = dimension_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()

    if data_source in always_include_data_source:
        run.text = data_source
    elif not_none(dimension):
        run.text = f"{fix_name(dimension)}: {fix_name(dim_label)}"
    else:
        return
        # run.text = "Overall"

    if len(run.text) > 33:
        run.text = run.text[:30] + '...'

    font = run.font
    font.name = 'Poppins'
    font.size = Pt(10)
    font.bold = True
    font.color.rgb = RGBColor(70, 177, 255)


def add_revenue_impact(anomaly_slide, left, top, width, height, revenue_impact, color, currency=None):
    shape = anomaly_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 176, 240)
    shadow = shape.shadow
    shadow.inherit = False
    line = shape.line
    line.fill.background()

    impact_textbox = anomaly_slide.shapes.add_textbox(
        Inches(left.inches + 0.05), top, width, height
    )

    impact_text_frame = impact_textbox.text_frame
    impact_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = impact_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()

    run.text = f"Rev Impact"
    font = run.font
    font.name = 'Poppins'
    font.size = Pt(9)
    font.color.rgb = RGBColor(255, 255, 255)

    p = impact_text_frame.add_paragraph()
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()

    run.text = f"{'-' if color in ['red1', 'red2'] else ''}{print_formatted(revenue_impact, 'Revenue', currency=currency)}"
    font = run.font
    font.name = 'Poppins'
    font.size = Pt(11) if len(run.text) < 8 else Pt(8)
    font.bold = True
    font.color.rgb = RGBColor(255, 255, 255)


def add_forecast_comment(anomaly_slide, left, top, width, height, metric, y, yhat, currency=None):
    forecast_comment_textbox = anomaly_slide.shapes.add_textbox(
        left, top, width, height
    )

    forecast_comment_text_frame = forecast_comment_textbox.text_frame
    forecast_comment_text_frame.word_wrap = True
    forecast_comment_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = forecast_comment_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()

    forecast_comment = f"{fix_name(metric)}"

    if metric.endswith('s'):
        forecast_comment = forecast_comment + " are"
    else:
        forecast_comment = forecast_comment + " is"

    forecast_comment = forecast_comment + f" {print_delta(now=y, prev=yhat)[1:]}"

    forecast_comment = forecast_comment + f" {'higher' if y > yhat else 'lower'} than expected value of {print_formatted(yhat, metric, currency=currency)}"

    run.text = forecast_comment
    font = run.font
    font.name = 'Poppins'
    font.size = Pt(10)
    font.bold = True
    font.color.rgb = RGBColor(255, 255, 255)


def add_divider(anomaly_slide, left, top, width, height, color):
    divider_shape = anomaly_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = divider_shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shadow = divider_shape.shadow
    shadow.inherit = False
    line = divider_shape.line
    line.fill.background()


def add_anomaly_chart(anomaly_slide, left, top, width, height, row):
    metric = row['metric']
    xaxis_data = row['xaxis_data']
    yaxis_data = row['yaxis_data']
    trend_data = row['trend_data']
    yhat_upper_data = row['yhat_upper_data']
    yhat_lower_data = row['yhat_lower_data']
    anomaly_type_data = row['yhat_anomaly_type_data']
    fig = get_graph(metric, xaxis_data, yaxis_data, trend_data, yhat_upper_data, yhat_lower_data, anomaly_type_data)
    img_path = "/tmp/anomaly.png" if in_production else "anomaly.png"
    fig.write_image(img_path, width=536, height=362)
    anomaly_slide.shapes.add_picture(
        img_path,
        left=left,
        top=top,
        width=width,
        height=height
    )


def add_footer(anomaly_slide, left, top, width, height, color):
    divider_shape = anomaly_slide.shapes.add_shape(
        MSO_SHAPE.ROUND_2_SAME_RECTANGLE, left, top, width, height
    )
    divider_shape.adjustments[0] = 0
    divider_shape.adjustments[1] = 0.1
    fill = divider_shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shadow = divider_shape.shadow
    shadow.inherit = False
    line = divider_shape.line
    line.fill.background()


def add_anomaly_card(anomaly_slide, row, period, left, top, color, currency=None):
    width = Cm(6.81)
    height = Cm(7.98)

    add_card_shape(
        anomaly_slide,
        left=left,
        top=top,
        width=width,
        height=height
    )

    if color == 'red':
        add_warning_critical(
            anomaly_slide,
            left=left,
            top=Inches(top.inches + 0),
            width=width,
            height=height,
            is_warning=row['is_warning'],
            is_critical=row['is_critical'],
        )

    add_metric_name(
        anomaly_slide,
        left=left,
        top=Inches(top.inches + 0.15),
        width=width,
        height=height,
        metric=row['metric']
    )

    add_metric_value(
        anomaly_slide,
        left=left,
        top=Inches(top.inches + 0.35),
        width=width,
        height=height,
        y=row['y'],
        metric=row['metric'],
        currency=currency
    )

    original_color = color

    if color == 'red' and row['metric'].lower():
        color = 'red1' if row['is_warning'] else 'red2'
    # else:
    #     color = 'green1' if row['is_warning'] else 'green2'
    if color == 'green' and row['metric'].lower():
        color = 'green1' if row['is_warning'] else 'green2'
    # else:
    #     color = 'red1' if row['is_warning'] else 'red2'
    if not (row['is_warning'] or row['is_critical']):
        color = 'grey'

    add_delta(
        anomaly_slide,
        left=Inches(left.inches + 1.06),
        top=Inches(top.inches + 0.38),
        width=width,
        height=height,
        period=period,
        y=row['y'],
        yhat=row['yhat'],
        color=color
    )

    add_revenue_impact(
        anomaly_slide,
        left=Inches(left.inches + 1.75),
        top=Inches(top.inches + 0.1),
        width=Cm(2.05),
        height=Cm(1.24),
        revenue_impact=row['revenue_impact'],
        color=color,
        currency=currency
    )

    add_dimension(
        anomaly_slide,
        left=left,
        top=Inches(top.inches + 0.65),
        width=width,
        height=height,
        dimension=row['dimension'],
        dim_label=row['dim_label'],
        data_source=row['data_source'],
    )

    add_divider(
        anomaly_slide,
        left=left,
        top=Inches(top.inches + 0.9),
        width=width,
        height=Cm(1.27),
        color=RGBColor(137, 219, 169) if original_color == 'green' else RGBColor(254, 180, 134)
    )

    add_forecast_comment(
        anomaly_slide,
        left=Inches(left.inches + width.inches / 2 - Cm(5.68).inches / 2),
        top=Inches(top.inches + 0.85),
        width=Cm(5.68),
        height=Cm(1.5),
        metric=row['metric'],
        y=row['y'],
        yhat=row['yhat'],
        currency=currency,
    )

    try:
        add_anomaly_chart(
            anomaly_slide,
            left=Inches(left.inches + width.inches / 2 - Cm(5.68).inches / 2),
            top=Inches(top.inches + 1.45),
            width=Cm(5.68),
            height=Cm(3.84),
            row=row
        )
    except Exception as e:
        print(f"No anomaly chart : {e}")

    # add_footer(
    #     anomaly_slide,
    #     left=left,
    #     top=Inches(top.inches + height.inches - Cm(1.28).inches),
    #     width=width,
    #     height=Cm(1.28),
    #     color=RGBColor(137, 219, 169) if color == 'green' else RGBColor(254, 180, 134)
    # )


def add_anomaly_cards(anomaly_slide, period, asset_df, kpi_list):

    top_row = Inches(0.75)
    bottom_row = Inches(4.25)

    # negative_card_locations = [
    #     (Inches(0.5), top_row),
    #     (Inches(3.5), top_row),
    #     (Inches(6.5), top_row),
    # ]
    #
    # positive_card_locations = [
    #     (Inches(0.5), bottom_row),
    #     (Inches(3.5), bottom_row),
    #     (Inches(6.5), bottom_row),
    # ]

    card_locations = [(Inches(0.5), top_row),
                      (Inches(3.5), top_row),
                      (Inches(6.5), top_row),
                      (Inches(0.5), bottom_row),
                      (Inches(3.5), bottom_row),
                      (Inches(6.5), bottom_row),
                      ]
    asset = asset_df.asset.unique()[0]

    # positive_anomaly_mask = asset_df['yhat_color'] == 'green'
    # negative_anomaly_mask = asset_df['yhat_color'] == 'red'
    if asset not in ['levis_indo_watchdog', 'levis_ph_watchdog', 'levis_ml_watchdog']:
        warning_or_critical_mask = (asset_df['is_warning'] == 1) | (asset_df['is_critical'] == 1)
    else:
        warning_or_critical_mask = ((asset_df['is_warning'] == 1) | (asset_df['is_critical'] == 1)) | \
                                   (asset_df['metric'].str.lower().isin(['units', 'orders', 'revenue', 'aur']))
    infinity_mask = np.isinf(asset_df['delta'])

    top_anomaly_df = asset_df[warning_or_critical_mask & (~infinity_mask)]
    top_anomaly_df['is_kpi'] = top_anomaly_df.apply(lambda row: is_kpi(row, asset_df, kpi_list), axis=1)
    top_anomaly_df['abs_revenue_impact'] = top_anomaly_df['revenue_impact'].apply(abs)
    top_anomaly_df.sort_values(by=['abs_revenue_impact'], ascending=False, inplace=True)

    # positive_anomaly_df = asset_df[positive_anomaly_mask & warning_or_critical_mask & (~infinity_mask)]
    # positive_anomaly_df['is_kpi'] = positive_anomaly_df.apply(lambda row: is_kpi(row, asset_df, kpi_list), axis=1)
    # positive_anomaly_df.sort_values(by=['is_kpi', 'revenue_impact'], ascending=False, inplace=True)
    #
    # negative_anomaly_df = asset_df[negative_anomaly_mask & warning_or_critical_mask & (~infinity_mask)]
    # negative_anomaly_df['is_kpi'] = negative_anomaly_df.apply(lambda row: is_kpi(row, asset_df, kpi_list), axis=1)
    # negative_anomaly_df.sort_values(by=['is_kpi', 'revenue_impact'], ascending=False, inplace=True)

    # max_metric = {}
    # for metric in top_anomaly_df.metric.unique():
    #     max_metric[metric] = top_anomaly_df[top_anomaly_df.metric == metric].revenue_impact.max()
    #
    # def impact_check(x):
    #     if max_metric[x['metric']] == 0:
    #         return False
    #     pct = ((max_metric[x['metric']] - x['revenue_impact']) * 100) / (max_metric[x['metric']])
    #     if pct < 15 and is_none(x['dimension']):
    #         return False
    #     else:
    #         return True
    #
    # top_anomaly_df['check'] = top_anomaly_df.apply(lambda x: impact_check(x), axis=1)

    # Check for currency of different clients
    currency = None
    for country in list(currency_dict.keys()):
        if f"_{country}_" in asset_df.asset.unique()[0]:
            currency = currency_dict[country]
            break
    top_anomaly_df = top_anomaly_df[top_anomaly_df['dimension'] != 'Device_Category']
    negative_warning_count = 0
    negative_critical_count = 0
    total_count = 0
    # & top_anomaly_df['check']
    for (i, row), (left, top) in zip(top_anomaly_df[(top_anomaly_df.y.notnull() & top_anomaly_df.y!=0) &
                                     ~top_anomaly_df.metric.str.contains('Return')].iterrows(), card_locations):
        if row['yhat_color'] == 'red':
            add_anomaly_card(anomaly_slide, row, period, left, top, color='red', currency=currency)
            total_count += 1
            if row['is_warning']:
                negative_warning_count += 1
            if row['is_critical']:
                negative_critical_count += 1
        else:
            add_anomaly_card(anomaly_slide, row, period, left, top, color='green', currency=currency)
            total_count += 1

    return negative_warning_count, negative_critical_count, total_count


def add_rca(slide, asset_df):

    print("adding RCA")

    rca_shape_left = Inches(9.15)
    rca_shape_top = Inches(1)
    rca_shape_width = Cm(6.81)
    rca_shape_height = Cm(7.98)

    dimension_textbox = slide.shapes.add_textbox(
        rca_shape_left, rca_shape_top, rca_shape_width, rca_shape_height
    )

    dimension_text_frame = dimension_textbox.text_frame
    dimension_text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = dimension_text_frame.paragraphs[0]
    print_revenue_rca(asset_df, p)


def add_anomaly_slide(ppt, period, asset, asset_df, kpi_list):
    blank_slide_layout = ppt.slide_layouts[6]
    anomaly_slide = ppt.slides.add_slide(blank_slide_layout)
    add_anomaly_heading(anomaly_slide, asset)

    # if asset != 'Levis SG':
    #     asset_df = asset_df[asset_df['data_source'] == 'capillary']

    if period == 'weekly':
        weekday = asset_df['weekday'].mode().iloc[0]
        if not asset_df[asset_df['weekday'] != weekday].empty:
            raise ValueError(f"Weekdays are different for {asset_df[asset_df['weekday'] != weekday]}")
    else:
        weekday = None

    add_date(anomaly_slide, period, weekday)
    # add_rca(anomaly_slide, asset_df)
    add_anomaly_legend(anomaly_slide)
    print(f"Creating {asset} anomaly cards...")
    negative_warning_count, negative_critical_count, total_count = add_anomaly_cards(anomaly_slide, period, asset_df, kpi_list)
    add_anomaly_subheadings(anomaly_slide, total_count)
    print(f"Created {asset} anomaly cards")

    return negative_warning_count, negative_critical_count, total_count
